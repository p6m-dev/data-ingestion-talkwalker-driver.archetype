import logging
import subprocess
from pptx import Presentation
from datetime import datetime
import os
import json
from . converter import Converter



class Pptx2Text(Converter):
    def __init__(self):
        Converter.__init__(self)        
        self.input_file_name = None
        self.output_file_name = None
        self.meta_output_file_name = None
        logging.basicConfig(
            format="%(asctime)s %(levelname)s: %(message)s", level=logging.DEBUG
        )
        self.logger = logging.getLogger()

    def configure(self, input_file_name: str, output_file_name: str):
        self.input_file_name = input_file_name
        self.output_file_name = output_file_name
        # Automatically generate meta_output_file_name based on output_file_name
        base_name, ext = os.path.splitext(output_file_name)
        self.meta_output_file_name = f"{base_name}_meta.json"

    def convert(self) -> bool:
        if self.input_file_name is None or self.output_file_name is None:
            self.logger.info("Pptx2Text.configure() call is missing.")
            return False

        try:
            self.logger.info(f"Converting PPTX from {self.input_file_name}.")

            # Check if the input file is in DOC format
            if self.input_file_name.lower().endswith('.ppt'):
                print(self.input_file_name)
                self.input_file_name = self.convert_ppt_to_pptx(self.input_file_name)


            presentation = Presentation(self.input_file_name)

            with open(self.output_file_name, 'w', encoding='utf-8') as txt_file:
                for slide in presentation.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            txt_file.write(shape.text + '\n')
            
            self.logger.info(f"{self.output_file_name} is written completely.")

            
            extracted_metadata = {}
            for attr in dir(presentation.core_properties):
                if not attr.startswith("_"):
                    try:
                        value = getattr(presentation.core_properties, attr)
                        if value is not None:
                            if isinstance(value, (str, int, float)):  # Accepts strings, integers, and floats
                                extracted_metadata[attr] = value
                            elif isinstance(value, datetime):
                                extracted_metadata[attr] = value.strftime('%Y-%m-%d %H:%M:%S')
                    except AttributeError:
                        pass

            print(extracted_metadata)

            # Save metadata in a JSON file
            json_metadata_file = self.meta_output_file_name
            with open(json_metadata_file, 'w', encoding='utf-8') as json_file:
                json.dump(extracted_metadata, json_file, indent=4)
            return True

        except Exception as e:
            logging.error(e)
            return False


    def convert_ppt_to_pptx(self, ppt_file_path):
            # Use subprocess to call the unoconv command
            subprocess.call(["unoconv", "-f", "pptx", ppt_file_path])

            return ppt_file_path.replace(".ppt", ".pptx")
    

def main():
    input_file = "data/pptxsample.pptx"  # Replace with your input PDF file path
    output_file = "data/pptxsample_output.txt"  # Replace with your desired output text file path
    
    pdf_converter = Pptx2Text()
    pdf_converter.configure(input_file, output_file)
    
    if pdf_converter.convert():
        print("Conversion successful!")
    else:
        print("Conversion failed.")

if __name__ == "__main__":
    main()